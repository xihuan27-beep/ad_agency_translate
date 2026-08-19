import io
import os
import subprocess
import tempfile
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_shapes(shapes):
    """Yield every leaf shape, recursing into groups."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _extract_from_text_frame(tf, s_idx, shape_id, prefix, text_units):
    """Append text units from a single text frame."""
    paragraphs = tf.paragraphs
    shape_lines = [p.text.strip() for p in paragraphs if len(p.text.strip()) >= 5]
    shape_text = "\n".join(shape_lines)
    for p_idx, para in enumerate(paragraphs):
        text = para.text.strip()
        if len(text) < 5:
            continue
        font_size = 14.0
        if para.runs and para.runs[0].font.size:
            font_size = para.runs[0].font.size.pt
        text_units.append({
            "id": f"s{s_idx}_sh{shape_id}_{prefix}_p{p_idx}",
            "slide_idx": s_idx,
            "shape_id": shape_id,
            "p_idx": p_idx,
            "ko_text": text,
            "font_size": font_size,
            "shape_text": shape_text,
            "shape_para_count": len(shape_lines),
        })


def extract_text_units(file_bytes: bytes) -> list[dict]:
    prs = Presentation(io.BytesIO(file_bytes))
    text_units = []
    seen_ids = set()

    for s_idx, slide in enumerate(prs.slides):
        for shape in _iter_shapes(slide.shapes):
            sid = shape.shape_id

            # Text frame (normal text boxes, placeholders)
            if shape.has_text_frame:
                _extract_from_text_frame(shape.text_frame, s_idx, sid, "tf", text_units)

            # Table cells
            elif shape.has_table:
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        _extract_from_text_frame(
                            cell.text_frame, s_idx, sid, f"r{r_idx}c{c_idx}", text_units
                        )

    # De-duplicate by id (merged cells repeat content)
    unique = []
    for u in text_units:
        if u["id"] not in seen_ids:
            seen_ids.add(u["id"])
            unique.append(u)
    return unique


def _apply_to_text_frame(tf, s_idx, shape_id, prefix, translations, font_name):
    orig_auto_size = tf.auto_size
    frame_modified = False
    for p_idx, para in enumerate(tf.paragraphs):
        unit_id = f"s{s_idx}_sh{shape_id}_{prefix}_p{p_idx}"
        if unit_id not in translations:
            continue
        en_text = translations[unit_id]
        ko_len = len(para.text)
        en_len = len(en_text)
        if ko_len == 0 or en_len == 0:
            continue
        run_sizes = [r.font.size.pt if r.font.size else None for r in para.runs]
        for run in para.runs:
            run.text = ""
        if not para.runs:
            continue
        para.runs[0].text = en_text
        if font_name:
            para.runs[0].font.name = font_name
        effective_ko_width = ko_len * 1.3
        if en_len > effective_ko_width:
            scale = max(effective_ko_width / en_len, 0.65)
            orig_size = run_sizes[0] if run_sizes else None
            if orig_size:
                para.runs[0].font.size = Pt(max(orig_size * scale, 8.0))
        frame_modified = True
    if frame_modified and orig_auto_size != MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass


def apply_translations(file_bytes: bytes, translations: dict[str, str], font_name: str = "") -> bytes:
    prs = Presentation(io.BytesIO(file_bytes))

    for s_idx, slide in enumerate(prs.slides):
        for shape in _iter_shapes(slide.shapes):
            sid = shape.shape_id
            if shape.has_text_frame:
                _apply_to_text_frame(shape.text_frame, s_idx, sid, "tf", translations, font_name)
            elif shape.has_table:
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        _apply_to_text_frame(
                            cell.text_frame, s_idx, sid, f"r{r_idx}c{c_idx}",
                            translations, font_name
                        )

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def render_slides_to_images(file_bytes: bytes) -> list[bytes]:
    """Convert each PPTX slide to a PNG image using LibreOffice headless.

    Returns a list of PNG bytes, one per slide (in order).
    Returns [] if LibreOffice is not available.
    """
    with tempfile.TemporaryDirectory() as tmpdir:
        pptx_path = os.path.join(tmpdir, "input.pptx")
        with open(pptx_path, "wb") as f:
            f.write(file_bytes)

        try:
            subprocess.run(
                [
                    "libreoffice", "--headless",
                    "--convert-to", "png",
                    "--outdir", tmpdir,
                    pptx_path,
                ],
                check=True,
                capture_output=True,
                timeout=120,
            )
        except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
            return []

        # LibreOffice names output files: input.png (1 slide) or input1.png input2.png ...
        pngs = sorted(
            [f for f in os.listdir(tmpdir) if f.endswith(".png")],
            key=lambda n: (len(n), n),  # lexicographic stable sort by length then name
        )
        result = []
        for name in pngs:
            with open(os.path.join(tmpdir, name), "rb") as f:
                result.append(f.read())
        return result


def extract_reference_texts(file_bytes: bytes) -> list[str]:
    """Extract all non-empty text lines from a reference (already-translated) PPTX."""
    prs = Presentation(io.BytesIO(file_bytes))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if len(text) >= 3:
                        texts.append(text)
    return texts
